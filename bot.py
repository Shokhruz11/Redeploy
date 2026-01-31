# -*- coding: utf-8 -*-
import os
import sqlite3
from datetime import datetime

import telebot
from telebot import types

from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document

# ============================
#        SOZLAMALAR
# ============================

BOT_TOKEN = os.getenv("BOT_TOKEN", "")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

ADMIN_TELEGRAM_ID = int(os.getenv("ADMIN_TELEGRAM_ID", "0"))  # raqam
ADMIN_USERNAME = os.getenv("ADMIN_USERNAME", "@admin")

CARD_NUMBER = os.getenv("CARD_NUMBER", "0000 0000 0000 0000")
CARD_OWNER = os.getenv("CARD_OWNER", "Karta egasi")

# Tariflar
FIRST_FREE = 1              # 1 ta topshiriq bepul
TASK_PRICE = 5000           # har bir keyingi topshiriq
REF_BONUS = 1000            # referal bonusi

GEMINI_MODEL = "gemini-1.5-flash"

# ============================
#        TEKSHIRUVLAR
# ============================

if not BOT_TOKEN:
    print("❌ BOT_TOKEN topilmadi. Railway Service Variables ga qo'shishni unutmang.")
if not GEMINI_API_KEY:
    print("❌ GEMINI_API_KEY topilmadi. AI ishlamaydi.")

# Telegram bot
bot = telebot.TeleBot(BOT_TOKEN)

# Gemini klienti
gemini_client = None
if GEMINI_API_KEY:
    try:
        gemini_client = genai.Client(api_key=GEMINI_API_KEY)
        print("✅ Gemini tayyor (google.genai).")
    except Exception as e:
        print("❌ Gemini ishga tushmadi:", e)

# ============================
#           DATABASE
# ============================

DB_PATH = "bot.db"


def db():
    conn = sqlite3.connect(DB_PATH, check_same_thread=False)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    conn = db()
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS users (
            user_id      INTEGER PRIMARY KEY,
            username     TEXT,
            first_name   TEXT,
            balance      INTEGER DEFAULT 0,
            free_used    INTEGER DEFAULT 0,
            referrals    INTEGER DEFAULT 0,
            created_at   TEXT
        );
        """
    )
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS payments (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id     INTEGER,
            amount      INTEGER,
            status      TEXT,
            created_at  TEXT
        );
        """
    )
    conn.commit()
    conn.close()


init_db()
print("✅ DB tayyor.")


def get_user(user):
    conn = db()
    cur = conn.cursor()
    cur.execute("SELECT * FROM users WHERE user_id = ?", (user.id,))
    row = cur.fetchone()
    if not row:
        cur.execute(
            """
            INSERT INTO users(user_id, username, first_name, created_at)
            VALUES(?,?,?,?)
            """,
            (user.id, user.username or "", user.first_name or "", datetime.utcnow().isoformat()),
        )
        conn.commit()
        cur.execute("SELECT * FROM users WHERE user_id = ?", (user.id,))
        row = cur.fetchone()
    conn.close()
    return row


def update_user_balance(user_id, delta):
    conn = db()
    cur = conn.cursor()
    cur.execute("UPDATE users SET balance = balance + ? WHERE user_id = ?", (delta, user_id))
    conn.commit()
    conn.close()


def mark_free_used(user_id):
    conn = db()
    cur = conn.cursor()
    cur.execute("UPDATE users SET free_used = 1 WHERE user_id = ?", (user_id,))
    conn.commit()
    conn.close()


def add_referral(ref_user_id):
    conn = db()
    cur = conn.cursor()
    cur.execute(
        "UPDATE users SET referrals = referrals + 1, balance = balance + ? WHERE user_id = ?",
        (REF_BONUS, ref_user_id),
    )
    conn.commit()
    conn.close()


# ============================
#            AI
# ============================

def ask_ai(prompt: str) -> str:
    if not gemini_client:
        return "ERROR: AI sozlanmagan. Admin GEMINI_API_KEY ni tekshirishi kerak."

    try:
        resp = gemini_client.models.generate_content(
            model=GEMINI_MODEL,
            contents=prompt,
        )
        text = (getattr(resp, "text", "") or "").strip()
        if not text:
            return "ERROR: AI javob qaytara olmadi."
        return text
    except Exception as e:
        print("Gemini xatosi:", repr(e))
        return "ERROR: AI xizmatida kutilmagan xatolik yuz berdi."


def build_prompt(kind: str, topic: str, pages: int) -> str:
    return f"""
Senga talabalar uchun yozma ish tayyorlash topshirig'i berilgan.

Ish turi: {kind}
Mavzu: {topic}
Taxminiy hajm: {pages} bet.

Talablar:
- Kirish, asosiy qism, xulosa bo'lsin.
- Ilmiy, adabiy va tushunarli uslubda yoz.
- Reja va sarlavhalarni ham kiriting.
- Matn o'zbek tilida bo'lsin.
- Keraksiz izohlar, "mana matn" kabi gaplar bo'lmasin.
"""


# ============================
#      FAYL GENERATSIYASI
# ============================

def create_pptx(title: str, text: str, filename: str) -> str:
    prs = Presentation()
    # Birinchi slayd – sarlavha
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = "AI tomonidan tayyorlangan taqdimot"

    # Qolgan slaydlar – matnni bo'lib joylashtiramiz
    slides_chunks = text.split("\n\n")
    for chunk in slides_chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for i, line in enumerate(chunk.split("\n")):
            if not line.strip():
                continue
            if i == 0:
                tf.text = line.strip()
            else:
                p = tf.add_paragraph()
                p.text = line.strip()
                p.level = 1
    prs.save(filename)
    return filename


def create_docx(title: str, text: str, filename: str) -> str:
    doc = Document()
    doc.add_heading(title, level=1)
    for par in text.split("\n\n"):
        doc.add_paragraph(par.strip())
    doc.save(filename)
    return filename


# ============================
#            MENYU
# ============================

def main_menu():
    kb = types.ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row("📝 Slayd", "📄 Mustaqil ish / Referat")
    kb.row("📚 Kurs ishi", "✍️ Insho", "🧩 Tezislar")
    kb.row("📰 Maqola", "🤖 AI bilan suhbat")
    kb.row("🎁 Referal bonus", "💰 Balans", "💳 To'lov / Hisob")
    kb.row("👨‍🏫 Profi jamoa", "❓ Yordam")
    return kb


def pay_menu():
    kb = types.InlineKeyboardMarkup()
    kb.add(
        types.InlineKeyboardButton("✅ 5 000 so'm", callback_data="pay_5000"),
        types.InlineKeyboardButton("✅ 10 000 so'm", callback_data="pay_10000"),
    )
    kb.add(types.InlineKeyboardButton("📸 Chekni yubordim", callback_data="pay_help"))
    return kb


# ============================
#         /start + REF
# ============================

@bot.message_handler(commands=["start"])
def cmd_start(message: types.Message):
    user = get_user(message.from_user)
    text = (
        "Assalomu alaykum, hurmatli talaba! 👋\n\n"
        "Siz *Super Talaba PRO* botidasiz.\n\n"
        "Bu yerda siz quyidagi xizmatlardan foydalanishingiz mumkin:\n"
        "• 📝 Slayd (PPTX)\n"
        "• 📄 Mustaqil ish / Referat (DOCX)\n"
        "• 📚 Kurs ishi\n"
        "• ✍️ Insho, 🧩 Tezislar, 📰 Maqola\n"
        "• 🤖 AI bilan suhbat\n"
        "• 🎁 Referal bonus, 💳 To'lov va chek tizimi\n\n"
        f"Bir dona topshiriq sizga *bepul* beriladi. Keyingi topshiriqlar narxi: *{TASK_PRICE} so'm*.\n\n"
        "Kerakli bo'limni menyudan tanlang."
    )

    # Referal: /start ref_123
    if message.text and " " in message.text:
        try:
            _, arg = message.text.split(" ", 1)
            if arg.startswith("ref_"):
                ref_id = int(arg.replace("ref_", ""))
                if ref_id != message.from_user.id:
                    add_referral(ref_id)
                    bot.send_message(
                        ref_id,
                        f"🎉 Do'stingiz referal havolangiz orqali kirdi. Sizga +{REF_BONUS} so'm bonus!",
                    )
        except Exception:
            pass

    bot.send_message(message.chat.id, text, reply_markup=main_menu(), parse_mode="Markdown")


# ============================
#      BALANS / REFERAL
# ============================

@bot.message_handler(commands=["balance"])
def cmd_balance(message: types.Message):
    user = get_user(message.from_user)
    ref_link = f"https://t.me/{bot.get_me().username}?start=ref_{message.from_user.id}"
    txt = (
        f"💰 Balansingiz: *{user['balance']} so'm*\n"
        f"🧾 Bepul ishlatilgan: {user['free_used']}/{FIRST_FREE}\n"
        f"👥 Referallar soni: {user['referrals']}\n\n"
        f"Referal havolangiz:\n`{ref_link}`"
    )
    bot.send_message(message.chat.id, txt, parse_mode="Markdown")


# ============================
#        TO'LOV / CHEK
# ============================

@bot.message_handler(commands=["pay"])
def cmd_pay(message: types.Message):
    text = (
        "💳 *To'lov / Hisobni to'ldirish*\n\n"
        f"Karta: `{CARD_NUMBER}`\n"
        f"Karta egasi: *{CARD_OWNER}*\n\n"
        "To'lov qilingandan so'ng chek (screenshot) ni shu botga yuboring.\n"
        "Admin chekni tekshirib, balansingizni to'ldirib beradi."
    )
    bot.send_message(message.chat.id, text, parse_mode="Markdown")


@bot.message_handler(content_types=["photo"])
def handle_cheque(message: types.Message):
    """Har qanday fotoni chek deb qabul qilamiz va admin-ga forward qilamiz."""
    user = get_user(message.from_user)
    caption = (
        f"🧾 Yangi chek!\n"
        f"👤 Foydalanuvchi: {message.from_user.first_name} (@{message.from_user.username})\n"
        f"🆔 ID: {message.from_user.id}\n"
        f"Balans: {user['balance']} so'm\n\n"
        "Quyidagi tugmalardan birini bosib tasdiqlang:"
    )

    kb = types.InlineKeyboardMarkup()
    kb.add(
        types.InlineKeyboardButton("✅ +5 000", callback_data=f"approve_{message.from_user.id}_5000"),
        types.InlineKeyboardButton("✅ +10 000", callback_data=f"approve_{message.from_user.id}_10000"),
    )
    kb.add(types.InlineKeyboardButton("❌ Rad etish", callback_data=f"approve_{message.from_user.id}_0"))

    try:
        bot.forward_message(ADMIN_TELEGRAM_ID, message.chat.id, message.message_id)
        bot.send_message(ADMIN_TELEGRAM_ID, caption, reply_markup=kb)
        bot.reply_to(message, "✅ Chek qabul qilindi. Admin tekshiradi, iltimos kuting.")
    except Exception as e:
        print("Chek forward xatosi:", e)
        bot.reply_to(message, "❌ Chekni adminga yuborishda xatolik yuz berdi.")


@bot.callback_query_handler(func=lambda c: c.data.startswith("approve_"))
def callback_approve(call: types.CallbackQuery):
    if call.from_user.id != ADMIN_TELEGRAM_ID:
        bot.answer_callback_query(call.id, "Faqat admin tasdiqlay oladi.")
        return

    _, user_id_str, amount_str = call.data.split("_")
    user_id = int(user_id_str)
    amount = int(amount_str)

    if amount > 0:
        update_user_balance(user_id, amount)
        conn = db()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO payments(user_id, amount, status, created_at) VALUES(?,?,?,?)",
            (user_id, amount, "approved", datetime.utcnow().isoformat()),
        )
        conn.commit()
        conn.close()
        bot.answer_callback_query(call.id, f"{amount} so'm tasdiqlandi.")
        bot.send_message(user_id, f"✅ To'lovingiz tasdiqlandi. Balansingiz +{amount} so'm.")
    else:
        bot.answer_callback_query(call.id, "Chek rad etildi.")
        bot.send_message(user_id, "❌ Kechirasiz, chek rad etildi. Admin bilan bog'laning.")

    # Tugmalarni o'chirib qo'yamiz
    try:
        bot.edit_message_reply_markup(call.message.chat.id, call.message.message_id, reply_markup=None)
    except Exception:
        pass


# ============================
#         XIZMAT FUNKSIYA
# ============================

def can_use_service(user_id: int, chat_id: int) -> bool:
    conn = db()
    cur = conn.cursor()
    cur.execute("SELECT balance, free_used FROM users WHERE user_id = ?", (user_id,))
    row = cur.fetchone()
    conn.close()

    if row["free_used"] < FIRST_FREE:
        # bir marta bepul
        mark_free_used(user_id)
        bot.send_message(
            chat_id,
            f"🎁 Sizga *bir dona bepul* topshiriq berildi. Keyingi topshiriqlar narxi: {TASK_PRICE} so'm.",
            parse_mode="Markdown",
        )
        return True

    if row["balance"] < TASK_PRICE:
        bot.send_message(
            chat_id,
            f"❗ Balansingiz yetarli emas.\n"
            f"Joriy balans: *{row['balance']} so'm*\n"
            f"Bir topshiriq narxi: *{TASK_PRICE} so'm*\n\n"
            f"/pay buyrug'i yoki \"💳 To'lov / Hisob\" tugmasi orqali balansni to'ldiring.",
            parse_mode="Markdown",
        )
        return False

    update_user_balance(user_id, -TASK_PRICE)
    bot.send_message(
        chat_id,
        f"✅ Xizmat uchun *{TASK_PRICE} so'm* yechildi. Rahmat!\n"
        f"Yangi balansni /balance orqali ko'rishingiz mumkin.",
        parse_mode="Markdown",
    )
    return True


def ask_topic_and_pages(chat_id: int, kind_code: str, nice_name: str):
    msg = bot.send_message(
        chat_id,
        f"{nice_name} uchun mavzuni va taxminiy betlar sonini yuboring.\n\n"
        "Masalan: `O'z-o'zini anglash psixologiyasi, 8`",
    )
    bot.register_next_step_handler(msg, handle_topic_pages_step, kind_code, nice_name)


def parse_topic_pages(text: str):
    if "," in text:
        topic, pages_str = text.split(",", 1)
        topic = topic.strip()
        try:
            pages = int(pages_str.strip())
        except ValueError:
            pages = 8
    else:
        topic = text.strip()
        pages = 8
    if pages < 3:
        pages = 3
    if pages > 30:
        pages = 30
    return topic, pages


def handle_topic_pages_step(message: types.Message, kind_code: str, nice_name: str):
    user = get_user(message.from_user)
    if not can_use_service(user["user_id"], message.chat.id):
        return

    topic, pages = parse_topic_pages(message.text)
    bot.send_message(message.chat.id, f"⌛ AI {nice_name} tayyorlamoqda. Kuting...")

    kind_text = {
        "slayd": "slayd taqdimot (PPTX)",
        "referat": "mustaqil ish / referat",
        "kurs": "kurs ishi",
        "insho": "insho",
        "tezis": "tezislar to'plami",
        "maqola": "ilmiy maqola",
    }.get(kind_code, "yozma ish")

    prompt = build_prompt(kind_text, topic, pages)
    ai_text = ask_ai(prompt)

    if ai_text.startswith("ERROR"):
        bot.send_message(message.chat.id, ai_text)
        return

    # Fayl yaratamiz
    try:
        if kind_code == "slayd":
            filename = f"slayd_{message.from_user.id}_{int(datetime.utcnow().timestamp())}.pptx"
            create_pptx(topic, ai_text, filename)
            with open(filename, "rb") as f:
                bot.send_document(
                    message.chat.id,
                    f,
                    visible_file_name=f"{topic[:30]} - slayd.pptx",
                    caption="✅ AI yordamida tayyorlangan slayd.",
                )
        else:
            filename = f"ish_{kind_code}_{message.from_user.id}_{int(datetime.utcnow().timestamp())}.docx"
            create_docx(topic, ai_text, filename)
            with open(filename, "rb") as f:
                bot.send_document(
                    message.chat.id,
                    f,
                    visible_file_name=f"{topic[:30]} - {kind_code}.docx",
                    caption=f"✅ AI yordamida tayyorlangan {kind_text}.",
                )
    except Exception as e:
        print("Fayl yaratish xatosi:", e)
        # Agar faylda muammo bo'lsa, matnni jo'natib yuboramiz
        bot.send_message(
            message.chat.id,
            "⚠️ Fayl yaratishda xatolik yuz berdi. Matnni shu yerga yuboryapman:\n\n" + ai_text,
        )


# ============================
#        MATNLI MENYU
# ============================

@bot.message_handler(func=lambda m: m.text in [
    "📝 Slayd",
    "📄 Mustaqil ish / Referat",
    "📚 Kurs ishi",
    "✍️ Insho",
    "🧩 Tezislar",
    "📰 Maqola",
    "🤖 AI bilan suhbat",
    "🎁 Referal bonus",
    "💰 Balans",
    "💳 To'lov / Hisob",
    "👨‍🏫 Profi jamoa",
    "❓ Yordam",
])
def menu_router(message: types.Message):
    text = message.text

    if text == "📝 Slayd":
        ask_topic_and_pages(message.chat.id, "slayd", "📝 Slayd (PPTX)")
    elif text == "📄 Mustaqil ish / Referat":
        ask_topic_and_pages(message.chat.id, "referat", "📄 Mustaqil ish / Referat")
    elif text == "📚 Kurs ishi":
        ask_topic_and_pages(message.chat.id, "kurs", "📚 Kurs ishi")
    elif text == "✍️ Insho":
        ask_topic_and_pages(message.chat.id, "insho", "✍️ Insho")
    elif text == "🧩 Tezislar":
        ask_topic_and_pages(message.chat.id, "tezis", "🧩 Tezislar")
    elif text == "📰 Maqola":
        ask_topic_and_pages(message.chat.id, "maqola", "📰 Maqola")
    elif text == "🤖 AI bilan suhbat":
        msg = bot.send_message(
            message.chat.id,
            "Savolingiz yoki mavzuni yozing. AI imkon qadar aniq va tushunarli javob beradi.",
        )
        bot.register_next_step_handler(msg, handle_chat_ai)
    elif text == "🎁 Referal bonus":
        user = get_user(message.from_user)
        ref_link = f"https://t.me/{bot.get_me().username}?start=ref_{message.from_user.id}"
        txt = (
            "🎁 *Referal dasturi*\n\n"
            "Do'stlaringizni taklif qiling va har bir aktiv foydalanuvchi uchun bonus oling!\n"
            f"Bonus: har bir do'stingiz uchun *{REF_BONUS} so'm*.\n\n"
            f"Referal havolasi:\n`{ref_link}`\n\n"
            f"Jami referallar: *{user['referrals']} ta*."
        )
        bot.send_message(message.chat.id, txt, parse_mode="Markdown")
    elif text == "💰 Balans":
        cmd_balance(message)
    elif text == "💳 To'lov / Hisob":
        cmd_pay(message)
    elif text == "👨‍🏫 Profi jamoa":
        bot.send_message(
            message.chat.id,
            "👨‍🏫 *Profi jamoa va katta hajmdagi ishlar*:\n\n"
            f"Kurs ishi, bitiruv malakaviy ishi yoki boshqa katta topshiriqlar uchun bevosita admin bilan bog'laning: {ADMIN_USERNAME}\n\n"
            "Shartlar va narxlar alohida kelishiladi.",
            parse_mode="Markdown",
        )
    elif text == "❓ Yordam":
        bot.send_message(
            message.chat.id,
            "❓ *Yordam bo'limi*\n\n"
            "1. /start – botni boshlash va menyuni ko'rish\n"
            "2. /balance – balansni ko'rish\n"
            "3. /pay – to'lov bo'limi\n\n"
            "Har qanday savol bo'lsa, bevosita admin bilan bog'lanishingiz mumkin: "
            f"{ADMIN_USERNAME}",
            parse_mode="Markdown",
        )


def handle_chat_ai(message: types.Message):
    user = get_user(message.from_user)
    if not can_use_service(user["user_id"], message.chat.id):
        return

    bot.send_message(message.chat.id, "🤖 AI javob tayyorlamoqda, kuting...")
    prompt = (
        "Talaba savoliga tushunarli, qisqa va aniq javob ber.\n\n"
        f"SAVOL/MAVZU:\n{message.text}"
    )
    answer = ask_ai(prompt)
    bot.send_message(message.chat.id, answer)


# ============================
#            RUN
# ============================

print("🚀 SUPER TALABA PRO bot ishga tushmoqda...")

if __name__ == "__main__":
    # long_polling parametrlarini biroz kichik qilish Railway uchun qulay
    bot.infinity_polling(timeout=10, long_polling_timeout=5)
